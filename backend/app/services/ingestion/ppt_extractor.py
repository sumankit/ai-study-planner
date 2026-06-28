# Validated in Colab Notebook 2, Cell 6
from pptx import Presentation
from typing import List, Dict

def extract_ppt(file_path: str) -> List[Dict]:
    prs    = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        combined = "\n".join(texts)
        if len(combined) > 20:
            slides.append({"page": i + 1, "text": combined})
    return slides